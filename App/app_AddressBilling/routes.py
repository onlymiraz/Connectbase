# /App/app_AddressBilling/routes.py
"""
==============================================================================
 HOLISTIC WORKFLOW FOR ADDRESS BILLING (High-Level Overview)
 ----------------------------------------------------------------------------
 1) A user uploads an address file (Excel/CSV/TXT) via the Flask front end.
    - The upload route: /upload => Renders upload.html
    - The user selects a file, enters an email, hits Submit.

 2) If CSV/TXT => parse => mapping => insertion into [ADDRESS_BILLING].[UI_LZ].
    If Excel => pick sheet => mapping => insertion => process_status='pending'.

 3) A separate scheduled ETL (fuzzymatch_script.py) picks up
    UI_LZ rows with 'pending' => matches => writes Fuzzymatch_Output => 'done'.

 4) The user can see results at /batch_history or
    /show_fuzzymatch_results/<batch_id>, plus an email on completion.

 5) Eventually, we archive 'done' rows => UI_LZ_Archive + Fuzzymatch_Output_Archive.

 6) The front-end pages:
    - upload.html, select_sheet.html, mapping.html, batch_history.html,
      fuzzymatch_results.html, etc.
==============================================================================
"""

import os
import getpass
import uuid
import csv
import pyodbc
import logging
import warnings
import psutil
import time
from functools import wraps
from logging.handlers import RotatingFileHandler

from flask import (
    request, render_template, redirect, url_for, flash,
    session, current_app, send_file, Response
)
from werkzeug.utils import secure_filename
from datetime import datetime

from app_AddressBilling.logging_config import logger
from app_AddressBilling.orchestration.ETL.odbc import _get_server_and_db, odbc_read, odbc_write

warnings.simplefilter('ignore', category=FutureWarning)

##############################################################################
# 1) Performance Logger Setup
##############################################################################
LOG_DEBUG_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'debug')
os.makedirs(LOG_DEBUG_DIR, exist_ok=True)
PERFORMANCE_LOG = os.path.join(LOG_DEBUG_DIR, 'performance.log')
performance_logger = logging.getLogger('performance_logger')
performance_logger.setLevel(logging.DEBUG)
if not performance_logger.handlers:
    perf_fh = RotatingFileHandler(
        PERFORMANCE_LOG,
        maxBytes=5 * 1024 * 1024,
        backupCount=5,
        encoding='utf-8'
    )
    perf_fh.setLevel(logging.DEBUG)
    perf_fmt = logging.Formatter(
        '%(asctime)s | PERF | %(funcName)s:%(lineno)d | %(message)s',
        datefmt='%Y-%m-%d %H:%M:%S'
    )
    perf_fh.setFormatter(perf_fmt)
    performance_logger.addHandler(perf_fh)

def log_performance(route_function):
    @wraps(route_function)
    def wrapper(*args, **kwargs):
        start_time = time.time()
        process = psutil.Process()
        mem_before = process.memory_info().rss
        cpu_before = psutil.cpu_percent(interval=None)

        path = request.path
        method = request.method
        user = request.remote_user or getpass.getuser()

        result = route_function(*args, **kwargs)

        elapsed = time.time() - start_time
        mem_after = process.memory_info().rss
        cpu_after = psutil.cpu_percent(interval=None)
        performance_logger.debug(
            f"PATH={path} | METHOD={method} | USER={user} | "
            f"Elapsed={elapsed:.3f}s | MemBefore={mem_before} | MemAfter={mem_after} | "
            f"CPUBefore={cpu_before} | CPUAfter={cpu_after}"
        )
        return result
    return wrapper

##############################################################################
# 2) CSV Helpers + Email
##############################################################################
def detect_delimiter(file_content):
    import csv
    try:
        dialect = csv.Sniffer().sniff(file_content[:1024], delimiters=',;|\t')
        return dialect.delimiter
    except:
        for delim in [',',';','|','\t']:
            if delim in file_content:
                return delim
        return ','

def read_csv_autodetect(file_path):
    import pandas as pd
    import csv as csv_module

    encodings_to_try = ['utf-8-sig','utf-8','latin1','iso-8859-1','cp1252']
    df = None
    for enc in encodings_to_try:
        try:
            with open(file_path, 'r', encoding=enc) as f:
                sample = f.read(2048)
            guessed_delim = detect_delimiter(sample)
            logger.info(f"Trying enc='{enc}', guessed_delim='{guessed_delim}'")

            read_args = {
                'filepath_or_buffer': file_path,
                'encoding': enc,
                'delimiter': guessed_delim,
                'on_bad_lines': 'skip',
                'dtype': str,
                'header': 0
            }
            if guessed_delim == '\t':
                read_args.update({
                    'engine': 'python',
                    'quotechar': None,
                    'quoting': csv_module.QUOTE_NONE
                })

            df_temp = pd.read_csv(**read_args)
            if len(df_temp.columns) == 1:
                forced_args = {
                    'filepath_or_buffer': file_path,
                    'encoding': enc,
                    'sep': '\t',
                    'engine': 'python',
                    'quotechar': None,
                    'quoting': csv_module.QUOTE_NONE,
                    'on_bad_lines': 'skip',
                    'dtype': str,
                    'header': 0
                }
                df_temp2 = pd.read_csv(**forced_args)
                if len(df_temp2.columns) > 1:
                    df_temp = df_temp2

            df = df_temp
            break
        except Exception as e:
            logger.warning(f"read_csv_autodetect fails with enc='{enc}': {e}")
    return df

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in {'xls','xlsx','csv','txt'}

def get_base_url_for_email():
    import socket
    hostname = socket.gethostname().upper()
    if hostname.endswith('DDV01'):
        return "https://stg-wad.ftr.com"
    elif hostname.endswith('DPV02') or hostname.endswith('DPV01'):
        return "https://wad.ftr.com"
    return "https://wad.ftr.com"

def send_upload_confirmation_email(batch_id, user_email):
    if not user_email:
        logger.info(f"No user_email for batch_id={batch_id}. Skipping upload email.")
        return
    base_url = get_base_url_for_email()
    link = f"{base_url}/show_fuzzymatch_results/{batch_id}"
    subject = f"Upload Confirmation for Batch {batch_id}"
    body = (
        f"Hello,\n\n"
        f"Your file has been successfully uploaded and is being processed.\n"
        f"The batch ID is {batch_id}.\n"
        f"Within about 2 hours, you'll receive another email with a link to the final matched data.\n"
        f"In the meantime, check /batch_history for progress.\n\n"
        f"Regards,\nWholesale Team"
    )
    from_addr = "WAD@ftr.com"
    to_addr = user_email

    import smtplib
    from email.mime.text import MIMEText
    from email.mime.multipart import MIMEMultipart
    msg = MIMEMultipart()
    msg["Subject"] = subject
    msg["From"] = from_addr
    msg["To"] = to_addr
    msg.attach(MIMEText(body, "plain"))

    try:
        with smtplib.SMTP("MailRelay.corp.pvt", 25) as smtp:
            smtp.send_message(msg)
        logger.info(f"Sent upload confirmation email to {user_email} for batch_id={batch_id}")
    except Exception as ex:
        logger.error(f"Failed sending upload email for batch_id={batch_id}, user_email={user_email}: {ex}")

def send_fuzzymatch_completion_email(batch_id, user_email):
    if not user_email:
        logger.info(f"No user_email found for batch_id={batch_id}, skipping completion email.")
        return
    base_url = get_base_url_for_email()
    link = f"{base_url}/show_fuzzymatch_results/{batch_id}"
    subject = f"Fuzzymatch Complete for Batch {batch_id}"
    body = (
        f"Hello,\n\n"
        f"Your fuzzymatch + billing data for batch {batch_id} is ready.\n"
        f"You can view it here:\n{link}\n\n"
        f"Regards,\nWholesale Team"
    )
    from_addr = "WAD@ftr.com"
    to_addr = user_email

    import smtplib
    from email.mime.text import MIMEText
    from email.mime.multipart import MIMEMultipart
    msg = MIMEMultipart()
    msg["Subject"] = subject
    msg["From"] = from_addr
    msg["To"] = to_addr
    msg.attach(MIMEText(body, "plain"))

    try:
        with smtplib.SMTP("MailRelay.corp.pvt", 25) as smtp:
            smtp.send_message(msg)
        logger.info(f"Sent fuzzymatch completion email to {user_email} for batch_id={batch_id}")
    except Exception as ex:
        logger.error(f"Failed sending fuzzymatch completion email for {batch_id}, user_email={user_email}: {ex}")

##############################################################################
# 3) Upload Flow: /upload, /select_sheet, /mapping, /process_mapping
##############################################################################
@log_performance
def upload():
    import pandas as pd
    from flask import redirect, request, flash, session, render_template

    if request.method == 'POST':
        user_email = request.form.get('user_email','').strip()
        if not user_email:
            flash("Please provide an email address for your results.")
            return redirect(request.url)

        session['user_email'] = user_email
        if 'file' not in request.files:
            flash('No file part in request.')
            return redirect(request.url)

        file = request.files['file']
        if file.filename == '':
            flash('No selected file.')
            return redirect(request.url)

        if file and allowed_file(file.filename):
            filename = secure_filename(file.filename)
            upload_folder = os.path.join(
                os.path.dirname(os.path.abspath(__file__)), 'uploads'
            )
            os.makedirs(upload_folder, exist_ok=True)
            file_path = os.path.join(upload_folder, filename)
            file.save(file_path)
            session['file_path'] = file_path

            try:
                if filename.lower().endswith(('.csv','.txt')):
                    df_csv = read_csv_autodetect(file_path)
                    if df_csv is None or df_csv.empty:
                        flash("Unable to parse CSV/TXT. Possibly bad encoding/delimiter.")
                        return redirect(request.url)

                    server, db = _get_server_and_db()
                    conn_str = (
                        f"DRIVER={{SQL Server}};"
                        f"SERVER={server};"
                        f"DATABASE={db};"
                        "Trusted_Connection=yes;"
                    )
                    conn = pyodbc.connect(conn_str)
                    cur = conn.cursor()
                    cur.execute("""
                        SELECT COLUMN_NAME, IS_NULLABLE, DATA_TYPE, CHARACTER_MAXIMUM_LENGTH
                        FROM INFORMATION_SCHEMA.COLUMNS
                        WHERE TABLE_NAME='UI_LZ'
                          AND TABLE_SCHEMA='ADDRESS_BILLING'
                          AND COLUMN_NAME NOT IN (
                            'ID','DtmStamp','ingestion_timestamp','batch_id',
                            'user_name','user_corp'
                          )
                        ORDER BY ORDINAL_POSITION
                    """)
                    sql_columns = {}
                    for row in cur.fetchall():
                        col_name = row.COLUMN_NAME
                        sql_columns[col_name] = {
                            'nullable': row.IS_NULLABLE,
                            'data_type': row.DATA_TYPE,
                            'max_length': row.CHARACTER_MAXIMUM_LENGTH
                        }
                    conn.close()

                    return render_template(
                        'app_AddressBilling/mapping.html',
                        file_columns=df_csv.columns.tolist(),
                        sql_columns=sql_columns,
                        sheet_name=None
                    )
                else:
                    excel_file = pd.ExcelFile(file_path)
                    sheet_names = excel_file.sheet_names
                    return render_template(
                        'app_AddressBilling/select_sheet.html',
                        sheet_names=sheet_names
                    )

            except Exception as e:
                logger.error(f"Error reading file: {e}")
                flash(f"Error reading file: {e}")
                return redirect(request.url)

    # GET => show form
    return render_template('app_AddressBilling/upload.html')

@log_performance
def select_sheet():
    from flask import flash, request, redirect, url_for, render_template, session
    import pandas as pd
    import os

    file_path = session.get('file_path')
    if not file_path or not os.path.exists(file_path):
        flash('File not found. Please upload again.')
        return redirect(url_for('upload'))

    if request.method == 'POST':
        sheet_name = request.form.get('sheet_name')
        if not sheet_name:
            flash("No sheet selected. Please pick a sheet.")
            return redirect(url_for('upload'))

        session['sheet_name'] = sheet_name
        return redirect(url_for('mapping'))

    try:
        excel_file = pd.ExcelFile(file_path)
        sheet_names = excel_file.sheet_names
        return render_template('app_AddressBilling/select_sheet.html', sheet_names=sheet_names)
    except Exception as e:
        logger.error(f"Error reading Excel file: {e}")
        flash(f"Error reading Excel file: {e}")
        return redirect(url_for('upload'))

@log_performance
def mapping():
    """
    Now allows GET or POST. If GET => possibly re-check session or show fallback.
    """
    from flask import flash, request, redirect, url_for, render_template, session
    import pandas as pd
    import os

    # If GET => we can either show a simple message or replicate logic:
    if request.method == 'GET':
        # Maybe user typed /mapping directly. Check if we have a file path and sheet name, etc.
        file_path = session.get('file_path')
        if not file_path or not os.path.exists(file_path):
            flash('No file path in session. Please upload again.')
            return redirect(url_for('upload'))
        # We can re-run the same logic or just show a friendly message:
        # For simplicity, we re-run the below logic. Then if it’s Excel, we read from session. 
        # If CSV => parse again. 
        # If no sheet => they must do /select_sheet.

    try:
        file_path = session.get('file_path')
        if not file_path or not os.path.exists(file_path):
            flash('No file path. Please upload again.')
            return redirect(url_for('upload'))

        # Decide CSV or Excel
        if file_path.lower().endswith(('.csv','.txt')):
            df = read_csv_autodetect(file_path)
            if df is None or df.empty:
                flash("Unable to parse CSV/TXT.")
                return redirect(url_for('upload'))
            sheet_name = None
        else:
            sheet_name = session.get('sheet_name')
            if not sheet_name:
                flash("No sheet_name stored in session. Please select a sheet first.")
                return redirect(url_for('upload'))

            df_obj = pd.read_excel(file_path, sheet_name=sheet_name)
            if isinstance(df_obj, dict):
                first_sheet = list(df_obj.keys())[0]
                df = df_obj[first_sheet]
                flash("Multiple DataFrames found; using the first by default.")
            else:
                df = df_obj

            if df.empty:
                flash(f"Sheet '{sheet_name}' is empty or unreadable.")
                return redirect(url_for('upload'))

        file_columns = df.columns.tolist()

        # query UI_LZ from STG/PRD
        server, db = _get_server_and_db()
        conn_str = (
            f"DRIVER={{SQL Server}};"
            f"SERVER={server};"
            f"DATABASE={db};"
            "Trusted_Connection=yes;"
        )
        conn = pyodbc.connect(conn_str)
        cur = conn.cursor()
        cur.execute("""
            SELECT COLUMN_NAME, IS_NULLABLE, DATA_TYPE, CHARACTER_MAXIMUM_LENGTH
            FROM INFORMATION_SCHEMA.COLUMNS
            WHERE TABLE_NAME='UI_LZ'
              AND TABLE_SCHEMA='ADDRESS_BILLING'
              AND COLUMN_NAME NOT IN (
                'ID','DtmStamp','ingestion_timestamp','batch_id',
                'user_name','user_corp'
              )
            ORDER BY ORDINAL_POSITION
        """)
        sql_columns = {}
        for row in cur.fetchall():
            col_name = row.COLUMN_NAME
            sql_columns[col_name] = {
                'nullable': row.IS_NULLABLE,
                'data_type': row.DATA_TYPE,
                'max_length': row.CHARACTER_MAXIMUM_LENGTH
            }
        conn.close()

        return render_template(
            'app_AddressBilling/mapping.html',
            file_columns=file_columns,
            sql_columns=sql_columns,
            sheet_name=sheet_name
        )
    except Exception as e:
        logger.error(f"Error in mapping: {e}")
        flash(f"Error in mapping: {e}")
        return redirect(url_for('upload'))

@log_performance
def process_mapping():
    from flask import flash, request, redirect, url_for, session
    import pandas as pd
    import os

    try:
        form_sheet_name = request.form.get('sheet_name','')
        session_sheet_name = session.get('sheet_name','')
        final_sheet_name = form_sheet_name if form_sheet_name else session_sheet_name

        mappings = {}
        for k,v in request.form.items():
            if k not in ('sheet_name',) and v.strip():
                mappings[k] = v

        file_path = session.get('file_path')
        if not file_path or not os.path.exists(file_path):
            flash('File not found. Please upload again.')
            return redirect(url_for('upload'))

        if file_path.lower().endswith(('.csv','.txt')):
            df = read_csv_autodetect(file_path)
            if df is None or df.empty:
                flash("Unable to parse CSV/TXT.")
                return redirect(url_for('upload'))
        else:
            if not final_sheet_name:
                flash("No sheet_name found. Please select a sheet first.")
                return redirect(url_for('upload'))

            df_obj = pd.read_excel(file_path, sheet_name=final_sheet_name)
            if isinstance(df_obj, dict):
                first_sheet = list(df_obj.keys())[0]
                df = df_obj[first_sheet]
                flash("Multiple DataFrames => used the first sheet automatically.")
            else:
                df = df_obj

            if df.empty:
                flash(f"Sheet '{final_sheet_name}' is empty or unreadable.")
                return redirect(url_for('upload'))

        df_mapped = pd.DataFrame()
        for sql_col, file_col in mappings.items():
            if file_col not in df.columns:
                flash(f"Column '{file_col}' not found in the file.")
                return redirect(url_for('upload'))
            df_mapped[sql_col] = df[file_col].astype(str)

        batch_id = str(uuid.uuid4())
        df_mapped['batch_id'] = batch_id
        df_mapped['user_name'] = getpass.getuser()
        df_mapped['user_corp'] = os.environ.get('USERDOMAIN','Unknown')
        df_mapped['ingestion_timestamp'] = pd.Timestamp.now()
        user_email = session.get('user_email','')
        df_mapped['user_email'] = user_email

        server, db = _get_server_and_db()
        conn_str = (
            f"DRIVER={{SQL Server}};"
            f"SERVER={server};"
            f"DATABASE={db};"
            "Trusted_Connection=yes;"
        )
        conn = pyodbc.connect(conn_str)
        cur = conn.cursor()

        cur.execute("SELECT COUNT(*) FROM [ADDRESS_BILLING].[UI_LZ]")
        before_count = cur.fetchval()

        columns = ', '.join(df_mapped.columns)
        placeholders = ', '.join(['?'] * len(df_mapped.columns))
        sql_insert = f"""
            INSERT INTO [ADDRESS_BILLING].[UI_LZ] ({columns})
            VALUES ({placeholders})
        """
        cur.fast_executemany = True
        cur.executemany(sql_insert, df_mapped.values.tolist())
        conn.commit()

        cur.execute("SELECT COUNT(*) FROM [ADDRESS_BILLING].[UI_LZ]")
        after_count = cur.fetchval()
        inserted_rows = after_count - before_count
        flash(f"Batch ID: {batch_id}, Rows Inserted: {inserted_rows}. We'll email you once processed!")

        cur.close()
        conn.close()

        send_upload_confirmation_email(batch_id, user_email)

        session.pop('file_path', None)
        session.pop('sheet_name', None)
        return redirect(url_for('upload'))

    except Exception as e:
        logger.error(f"Error in process_mapping: {e}")
        flash(f"Error in process_mapping: {e}")
        return redirect(url_for('upload'))

@log_performance
def get_batch_data():
    from flask import flash, request, redirect, url_for
    import pandas as pd
    from io import StringIO, BytesIO

    if request.method == 'POST':
        batch_id = request.form.get('batch_id','')
        user_name = getpass.getuser()
        try:
            server, db = _get_server_and_db()
            conn_str = (
                f"DRIVER={{SQL Server}};"
                f"SERVER={server};"
                f"DATABASE={db};"
                "Trusted_Connection=yes;"
            )
            conn = pyodbc.connect(conn_str)
            query = """
                SELECT *
                FROM [ADDRESS_BILLING].[UI_LZ]
                WHERE batch_id=? AND user_name=?
            """
            df = pd.read_sql(query, conn, params=(batch_id, user_name))
            conn.close()
            if df.empty:
                flash(f"No data found for batch {batch_id} & user {user_name}")
                return redirect(url_for('get_batch_data'))

            output = StringIO()
            df.to_csv(output, index=False)
            output.seek(0)
            return send_file(
                BytesIO(output.getvalue().encode()),
                mimetype='text/csv',
                download_name=f'batch_{batch_id}_data.csv'
            )
        except Exception as e:
            logger.error(f"Error retrieving batch data: {e}")
            flash(f"Error retrieving batch data: {e}")
            return redirect(url_for('get_batch_data'))

    return render_template('app_AddressBilling/get_batch_data.html')

##############################################################################
# 4) fuzzymatch_powerbi_view, show_fuzzymatch_results, summary, downloads
##############################################################################
@log_performance
def fuzzymatch_powerbi_view(batch_id):
    import pandas as pd
    from flask import flash, redirect, url_for, render_template

    try:
        server, db = _get_server_and_db()
        conn_str = (
            f"DRIVER={{SQL Server}};"
            f"SERVER={server};"
            f"DATABASE={db};"
            "Trusted_Connection=yes;"
        )
        conn = pyodbc.connect(conn_str)
        query = "SELECT * FROM [ADDRESS_BILLING].[Fuzzymatch_Output] WHERE batch_id=?"
        df = pd.read_sql(query, conn, params=[batch_id])
        conn.close()

        if df.empty:
            flash(f"No fuzzymatch results for batch {batch_id}.")
            return redirect(url_for('upload'))

        columns = df.columns.tolist()
        data = df.values.tolist()
        return render_template(
            'app_AddressBilling/fuzzymatch_powerbi_view.html',
            batch_id=batch_id,
            columns=columns,
            data=data
        )
    except Exception as e:
        logger.error(f"Error loading fuzzymatch_powerbi_view: {e}")
        flash(f"Error loading fuzzymatch_powerbi_view: {e}")
        return redirect(url_for('upload'))

@log_performance
def show_fuzzymatch_results(batch_id):
    import pandas as pd
    from flask import flash, redirect, url_for, render_template

    MAX_DISPLAY = 10000
    try:
        server, db = _get_server_and_db()
        conn_str = (
            f"DRIVER={{SQL Server}};"
            f"SERVER={server};"
            f"DATABASE={db};"
            "Trusted_Connection=yes;"
        )
        conn = pyodbc.connect(conn_str)

        df_cnt_main = pd.read_sql(
            "SELECT COUNT(*) AS cnt FROM [ADDRESS_BILLING].[Fuzzymatch_Output] WHERE batch_id=?",
            conn, params=[batch_id]
        )
        main_count = df_cnt_main['cnt'].iloc[0]

        if main_count == 0:
            df_cnt_arch = pd.read_sql(
                "SELECT COUNT(*) AS cnt FROM [ADDRESS_BILLING].[Fuzzymatch_Output_Archive] WHERE batch_id=?",
                conn, params=[batch_id]
            )
            arch_count = df_cnt_arch['cnt'].iloc[0]
            if arch_count == 0:
                df_lz = pd.read_sql(
                    "SELECT * FROM [ADDRESS_BILLING].[UI_LZ] WHERE batch_id=?",
                    conn, params=[batch_id]
                )
                if df_lz.empty:
                    df_lz_arch = pd.read_sql(
                        "SELECT * FROM [ADDRESS_BILLING].[UI_LZ_Archive] WHERE batch_id=?",
                        conn, params=[batch_id]
                    )
                    if df_lz_arch.empty:
                        conn.close()
                        flash(f"No fuzzymatch data nor raw data for batch {batch_id}.")
                        return redirect(url_for('upload'))
                    else:
                        df_lz = df_lz_arch

                conn.close()
                columns_list = df_lz.columns.tolist()
                data_rows = df_lz.values.tolist()
                notice = "No fuzzymatch rows found. Showing raw LZ data instead."
                return render_template(
                    'app_AddressBilling/fuzzymatch_results.html',
                    batch_id=batch_id,
                    columns_list=columns_list,
                    data_rows=data_rows,
                    total_rows=len(df_lz),
                    notice=notice
                )

            top_clause = f"TOP {MAX_DISPLAY}" if arch_count > MAX_DISPLAY else ""
            sql_arch = f"""
                SELECT {top_clause} *
                FROM [ADDRESS_BILLING].[Fuzzymatch_Output_Archive]
                WHERE batch_id=?
                ORDER BY ID
            """
            df = pd.read_sql(sql_arch, conn, params=[batch_id])
            total_rows = arch_count
        else:
            top_clause = f"TOP {MAX_DISPLAY}" if main_count > MAX_DISPLAY else ""
            sql_main = f"""
                SELECT {top_clause} *
                FROM [ADDRESS_BILLING].[Fuzzymatch_Output]
                WHERE batch_id=?
                ORDER BY ID
            """
            df = pd.read_sql(sql_main, conn, params=[batch_id])
            total_rows = main_count

        conn.close()

        if df.empty:
            flash(f"No fuzzymatch rows found for batch {batch_id}.")
            return redirect(url_for('upload'))

        notice = None
        if total_rows > MAX_DISPLAY:
            notice = f"Showing only first {MAX_DISPLAY} rows. Download to see all."

        columns_list = df.columns.tolist()
        data_rows = df.values.tolist()

        return render_template(
            'app_AddressBilling/fuzzymatch_results.html',
            batch_id=batch_id,
            columns_list=columns_list,
            data_rows=data_rows,
            total_rows=total_rows,
            notice=notice
        )

    except Exception as e:
        logger.error(f"Error in show_fuzzymatch_results({batch_id}): {e}")
        flash(f"Error showing fuzzymatch results: {e}")
        return redirect(url_for('upload'))

@log_performance
def show_fuzzymatch_summary(batch_id):
    import pandas as pd
    from flask import flash, redirect, url_for, render_template

    try:
        server, db = _get_server_and_db()
        conn_str = (
            f"DRIVER={{SQL Server}};"
            f"SERVER={server};"
            f"DATABASE={db};"
            "Trusted_Connection=yes;"
        )
        conn = pyodbc.connect(conn_str)

        main_count = pd.read_sql(
            "SELECT COUNT(*) AS cnt FROM [ADDRESS_BILLING].[Fuzzymatch_Output] WHERE batch_id=?",
            conn, params=[batch_id]
        )['cnt'].iloc[0]

        if main_count == 0:
            arch_count = pd.read_sql(
                "SELECT COUNT(*) AS cnt FROM [ADDRESS_BILLING].[Fuzzymatch_Output_Archive] WHERE batch_id=?",
                conn, params=[batch_id]
            )['cnt'].iloc[0]
            if arch_count == 0:
                conn.close()
                notice = "No fuzzymatch data found yet for this batch."
                pivot_html = "<p>(No matched data to summarize)</p>"
                return render_template(
                    'app_AddressBilling/fuzzymatch_summary.html',
                    batch_id=batch_id,
                    pivot_html=pivot_html,
                    notice=notice
                )
            else:
                df = pd.read_sql(
                    "SELECT * FROM [ADDRESS_BILLING].[Fuzzymatch_Output_Archive] WHERE batch_id=?",
                    conn, params=[batch_id]
                )
                conn.close()
        else:
            df = pd.read_sql(
                "SELECT * FROM [ADDRESS_BILLING].[Fuzzymatch_Output] WHERE batch_id=?",
                conn, params=[batch_id]
            )
            conn.close()

        if df.empty:
            notice = "No fuzzymatch data found for this batch."
            pivot_html = "<p>(No matched data to summarize)</p>"
            return render_template(
                'app_AddressBilling/fuzzymatch_summary.html',
                batch_id=batch_id,
                pivot_html=pivot_html,
                notice=notice
            )

        if 'Matched_State' not in df.columns:
            notice = "No 'Matched_State' column found."
            pivot_html = "<p>(No matched state data to pivot)</p>"
        else:
            pivot_df = df.pivot_table(
                index='Matched_State',
                values='Matched_Address',
                aggfunc='count'
            ).rename(columns={'Matched_Address':'Count_Addresses'}).reset_index()
            pivot_html = pivot_df.to_html(index=False, classes="table table-striped table-sm")
            notice = None

        return render_template(
            'app_AddressBilling/fuzzymatch_summary.html',
            batch_id=batch_id,
            pivot_html=pivot_html,
            notice=notice
        )
    except Exception as e:
        logger.error(f"Error in show_fuzzymatch_summary({batch_id}): {e}")
        flash(f"Summary error for {batch_id}: {e}")
        return redirect(url_for('upload'))

@log_performance
def download_fuzzymatch_csv(batch_id):
    from flask import flash, redirect, url_for, Response

    table_name, row_count = _which_fuzzymatch_table(batch_id)
    if not table_name or row_count == 0:
        flash(f"No data to download for batch {batch_id}")
        return redirect(url_for('show_fuzzymatch_results', batch_id=batch_id))

    def generate_csv():
        server, db = _get_server_and_db()
        conn_str = (
            f"DRIVER={{SQL Server}};"
            f"SERVER={server};"
            f"DATABASE={db};"
            "Trusted_Connection=yes;"
        )
        conn = pyodbc.connect(conn_str)
        cursor = conn.cursor()
        cursor.execute(f"SELECT * FROM {table_name} WHERE batch_id=? ORDER BY ID", (batch_id,))
        col_names = [desc[0] for desc in cursor.description]
        yield ",".join(col_names) + "\r\n"

        chunk = cursor.fetchmany(2000)
        while chunk:
            for row in chunk:
                cells = []
                for val in row:
                    if val is None:
                        cells.append('')
                    else:
                        s = str(val).replace('"','""')
                        if ',' in s or '\n' in s or '\r' in s:
                            s = f'"{s}"'
                        cells.append(s)
                yield ",".join(cells) + "\r\n"
            chunk = cursor.fetchmany(2000)

        cursor.close()
        conn.close()

    filename = f"Fuzzymatch_{batch_id}.csv"
    return Response(
        generate_csv(),
        mimetype='text/csv',
        headers={"Content-Disposition": f"attachment; filename={filename}"}
    )

@log_performance
def download_fuzzymatch_excel(batch_id):
    import openpyxl
    from openpyxl import Workbook
    from io import BytesIO
    from flask import flash, redirect, url_for, send_file

    table_name, row_count = _which_fuzzymatch_table(batch_id)
    if not table_name or row_count == 0:
        flash(f"No data to download for batch {batch_id}")
        return redirect(url_for('show_fuzzymatch_results', batch_id=batch_id))

    output = BytesIO()
    server, db = _get_server_and_db()
    conn_str = (
        f"DRIVER={{SQL Server}};"
        f"SERVER={server};"
        f"DATABASE={db};"
        "Trusted_Connection=yes;"
    )
    conn = pyodbc.connect(conn_str)
    cursor = conn.cursor()
    cursor.execute(f"SELECT * FROM {table_name} WHERE batch_id=? ORDER BY ID", (batch_id,))

    wb = openpyxl.Workbook(write_only=True)
    ws = wb.create_sheet("FuzzymatchData")
    if "Sheet" in wb.sheetnames:
        wb.remove(wb["Sheet"])

    col_names = [desc[0] for desc in cursor.description]
    ws.append(col_names)

    chunk = cursor.fetchmany(2000)
    while chunk:
        for row in chunk:
            row_list = [val if val is not None else "" for val in row]
            ws.append(row_list)
        chunk = cursor.fetchmany(2000)

    cursor.close()
    conn.close()

    wb.save(output)
    output.seek(0)
    filename = f"Fuzzymatch_{batch_id}.xlsx"
    return send_file(
        output,
        mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
        as_attachment=True,
        download_name=filename
    )

def _which_fuzzymatch_table(batch_id):
    server, db = _get_server_and_db()
    conn_str = (
        f"DRIVER={{SQL Server}};"
        f"SERVER={server};"
        f"DATABASE={db};"
        "Trusted_Connection=yes;"
    )
    conn = pyodbc.connect(conn_str)
    cursor = conn.cursor()

    cursor.execute("SELECT COUNT(*) FROM [ADDRESS_BILLING].[Fuzzymatch_Output] WHERE batch_id=?", (batch_id,))
    main_count = cursor.fetchone()[0]
    if main_count > 0:
        cursor.close()
        conn.close()
        return ("[ADDRESS_BILLING].[Fuzzymatch_Output]", main_count)

    cursor.execute("SELECT COUNT(*) FROM [ADDRESS_BILLING].[Fuzzymatch_Output_Archive] WHERE batch_id=?", (batch_id,))
    arch_count = cursor.fetchone()[0]
    cursor.close()
    conn.close()

    if arch_count > 0:
        return ("[ADDRESS_BILLING].[Fuzzymatch_Output_Archive]", arch_count)
    return (None, 0)

##############################################################################
# 5) batch_history => combine UI_LZ + UI_LZ_Archive with fuzzymatch counts
##############################################################################
@log_performance
def batch_history():
    import pandas as pd
    from flask import request, render_template, flash

    try:
        search_query = request.args.get('search_query','').strip()
        logger.info(f"batch_history() search_query='{search_query}'")

        server, db = _get_server_and_db()
        conn_str = (
            f"DRIVER={{SQL Server}};"
            f"SERVER={server};"
            f"DATABASE={db};"
            "Trusted_Connection=yes;"
        )
        conn = pyodbc.connect(conn_str)

        where_clause = ""
        params = []
        if search_query:
            like_str = f"%{search_query}%"
            where_clause = """
                WHERE
                  (batch_id LIKE ? OR
                   user_email LIKE ? OR
                   process_status LIKE ?)
            """
            params = [like_str, like_str, like_str]

        extra_condition = "WHERE rn=1" if not where_clause else "AND rn=1"

        base_cte = f"""
        WITH cte AS (
            SELECT
                batch_id,
                user_email,
                ingestion_timestamp,
                process_status,
                ROW_NUMBER() OVER (
                    PARTITION BY batch_id
                    ORDER BY ingestion_timestamp DESC
                ) AS rn
            FROM (
                SELECT batch_id, user_email, ingestion_timestamp, process_status
                FROM [ADDRESS_BILLING].[UI_LZ]
                UNION ALL
                SELECT batch_id, user_email, ingestion_timestamp, process_status
                FROM [ADDRESS_BILLING].[UI_LZ_Archive]
            ) u
            {where_clause}
        )
        SELECT
            cte.batch_id,
            cte.user_email,
            cte.ingestion_timestamp,
            cte.process_status
        FROM cte
        {extra_condition}
        ORDER BY cte.ingestion_timestamp DESC
        """
        df_core = pd.read_sql(base_cte, conn, params=params)
        if df_core.empty:
            conn.close()
            flash("No results found." if search_query else "No batches found.")
            return render_template('app_AddressBilling/batch_history.html', df=[])

        fuzz_sql = """
        SELECT x.batch_id, SUM(x.cnt) AS row_count
        FROM (
            SELECT batch_id, COUNT(*) AS cnt
            FROM [ADDRESS_BILLING].[Fuzzymatch_Output]
            GROUP BY batch_id
            UNION ALL
            SELECT batch_id, COUNT(*) AS cnt
            FROM [ADDRESS_BILLING].[Fuzzymatch_Output_Archive]
            GROUP BY batch_id
        ) x
        GROUP BY x.batch_id
        """
        df_fz = pd.read_sql(fuzz_sql, conn)
        conn.close()

        df_fz.rename(columns={'batch_id':'fz_batch','row_count':'fz_count'}, inplace=True)
        merged = pd.merge(df_core, df_fz, how='left', left_on='batch_id', right_on='fz_batch')
        merged['fz_count'] = merged['fz_count'].fillna(0).astype(int)

        final_records = merged.to_dict(orient='records')
        return render_template('app_AddressBilling/batch_history.html', df=final_records)
    except Exception as e:
        logger.error(f"batch_history error: {e}", exc_info=True)
        flash(f"Error loading batch history: {e}")
        return render_template('app_AddressBilling/batch_history.html', df=[])

##############################################################################
# 6) generate_documentation (PDF/Word/PPTX)
##############################################################################
from fpdf import FPDF
import io
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt

@log_performance
def generate_documentation(doc_type):
    from flask import flash, redirect, send_file, url_for

    doc_title = "Frontier Address Billing System"
    doc_subtitle = "Technical & Business Documentation"
    doc_overview = (
        "This system allows users to upload address files, map them to the "
        "database (UI_LZ), run fuzzy matching via an ETL, and store results in "
        "Fuzzymatch_Output. Data is eventually archived. The UI has a batch "
        "history page, results pages, and CSV/Excel downloads.\n"
    )

    doc_workflow = (
        "1) User uploads file via /upload\n"
        "2) The file is mapped => [ADDRESS_BILLING].[UI_LZ], process_status='pending'\n"
        "3) A scheduled ETL runs fuzzy matching => results in Fuzzymatch_Output\n"
        "4) The UI offers CSV/Excel + pivot summary\n"
        "5) ~2 hours, user gets an email with direct link\n"
        "6) Data is eventually archived to UI_LZ_Archive & Fuzzymatch_Output_Archive\n"
    )

    doc_sql = (
        "SQL Server Objects:\n"
        "- [ADDRESS_BILLING].[UI_LZ], [UI_LZ_Archive]\n"
        "- [ADDRESS_BILLING].[Fuzzymatch_Output], [Fuzzymatch_Output_Archive]\n"
        "- [ADDRESS_BILLING].[ADDR_BILLING_MASTER]\n"
        "- Index scripts, archive scripts, etc.\n"
    )

    doc_routes = (
        "Flask Endpoints:\n"
        "- /upload (CSV/TXT/Excel)\n"
        "- /select_sheet (Excel sheet selection)\n"
        "- /mapping => /process_mapping\n"
        "- /batch_history => list of all batches\n"
        "- /show_fuzzymatch_results/<batch_id>\n"
        "- /download_fuzzymatch_excel/<batch_id>, /download_fuzzymatch_csv/<batch_id>\n"
        "- /fuzzymatch_powerbi_view/<batch_id>\n"
        "- /show_fuzzymatch_summary/<batch_id>\n"
    )

    user_manual = (
        "USER MANUAL:\n\n"
        "1) Prepare addresses in Excel/CSV.\n"
        "2) /upload => fill out email => pick file => submit.\n"
        "3) If Excel => pick sheet => map columns to DB => done.\n"
        "4) Wait ~2 hours => get completion email.\n"
        "5) /batch_history => see your batches. Or use direct email link.\n"
        "6) Download data or see fuzzy results in the UI.\n"
    )

    project_structure = (
        "/\n"
        "  app.py, base.html, index.html...\n"
        "app_AddressBilling/\n"
        "  routes.py, select_sheet.html, mapping.html, etc.\n"
        "app_CircuitInventory/\n"
        "app_Locations/\n"
    )

    full_text = (
        f"{doc_title}\n\n{doc_subtitle}\n{'='*40}\n\n"
        f"[OVERVIEW]\n{doc_overview}\n"
        f"[WORKFLOW]\n{doc_workflow}\n"
        f"[SQL]\n{doc_sql}\n"
        f"[ROUTES]\n{doc_routes}\n"
        "----------------------------------\n"
        + user_manual + "\n"
        "----------------------------------\n\n"
        f"[PROJECT STRUCTURE]\n{project_structure}\n"
        "\nEnd of documentation.\n"
    )

    if doc_type == 'pdf':
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Arial", 'B', 14)
        pdf.cell(0, 10, doc_title, ln=1, align="C")
        pdf.set_font("Arial", '', 11)

        for line in full_text.split('\n'):
            safe_line = line.encode('ascii','ignore').decode('ascii','ignore')
            pdf.multi_cell(0, 6, safe_line)
        pdf_buffer = io.BytesIO(pdf.output(dest='S').encode('latin-1'))
        return send_file(
            pdf_buffer,
            mimetype='application/pdf',
            as_attachment=True,
            download_name="Address_Billing_UserManual.pdf"
        )

    elif doc_type == 'docx':
        doc = Document()
        doc.add_heading(doc_title, 0)
        p = doc.add_paragraph(doc_subtitle)
        p.style = 'Intense Quote'
        safe_text = ''.join(ch for ch in full_text if ord(ch) < 65535)
        doc.add_paragraph(safe_text)
        temp_stream = io.BytesIO()
        doc.save(temp_stream)
        temp_stream.seek(0)
        return send_file(
            temp_stream,
            mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            as_attachment=True,
            download_name="Address_Billing_UserManual.docx"
        )

    elif doc_type == 'pptx':
        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = doc_title
        slide.placeholders[1].text = doc_subtitle

        bullet_layout = prs.slide_layouts[1]
        # Slide #2 => overview + workflow
        slide2 = prs.slides.add_slide(bullet_layout)
        slide2.shapes.title.text = "Overview & Workflow"
        tf2 = slide2.shapes.placeholders[1].text_frame
        for l in (doc_overview + "\n" + doc_workflow).split('\n'):
            if l.strip():
                tf2.add_paragraph().text = l

        # Slide #3 => user manual
        slide3 = prs.slides.add_slide(bullet_layout)
        slide3.shapes.title.text = "User Manual"
        tf3 = slide3.shapes.placeholders[1].text_frame
        for l in user_manual.split('\n'):
            if l.strip():
                tf3.add_paragraph().text = l

        ppt_buffer = io.BytesIO()
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)
        return send_file(
            ppt_buffer,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            as_attachment=True,
            download_name="Address_Billing_UserManual.pptx"
        )
    else:
        flash(f"Unknown doc type: {doc_type}")
        return redirect(url_for('upload'))
